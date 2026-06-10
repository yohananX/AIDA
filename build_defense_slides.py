"""
Build AIDA defense slides matching Toluwanimi's reference format.
Widescreen 13.33x7.50 inches, Times New Roman, Title slide layout.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
FONT_NAME = "Times New Roman"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout (index 6 is typically blank in default template)
blank_layout = prs.slide_layouts[6]

def add_textbox(slide, left, top, width, height):
    """Add a textbox and return its text_frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    return txBox.text_frame

def set_para(tf, text, size=Pt(24), bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Append a paragraph to a text frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return p

def add_title_slide(title_text, subtitle_lines=None):
    """Create a slide with a centered title and optional subtitle lines."""
    slide = prs.slides.add_slide(blank_layout)
    # Title
    tf = add_textbox(slide, Inches(1), Inches(2.5), Inches(11.33), Inches(2))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.name = FONT_NAME
    run.font.size = Pt(40)
    run.font.bold = True
    # Subtitles
    if subtitle_lines:
        tf2 = add_textbox(slide, Inches(1), Inches(4.5), Inches(11.33), Inches(2))
        for line in subtitle_lines:
            set_para(tf2, line, size=Pt(20), bold=False, alignment=PP_ALIGN.CENTER)
    return slide

def add_section_slide(section_num, section_title):
    """Create a numbered section divider slide."""
    slide = prs.slides.add_slide(blank_layout)
    # Number
    tf = add_textbox(slide, Inches(0.5), Inches(0.3), Inches(2), Inches(1))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"{section_num:02d}"
    run.font.name = FONT_NAME
    run.font.size = Pt(24)
    run.font.bold = True
    # Title
    tf2 = add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.33), Inches(1))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = section_title
    run2.font.name = FONT_NAME
    run2.font.size = Pt(24)
    run2.font.bold = True
    # Separator line
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.5), Inches(2.1), Inches(12.33), Pt(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape.line.fill.background()
    return slide

def add_content_slide(title, bullets, title_size=Pt(24), bullet_size=Pt(20)):
    """Standard content slide with a title and bullet points."""
    slide = prs.slides.add_slide(blank_layout)
    # Title
    tf = add_textbox(slide, Inches(0.7), Inches(0.3), Inches(12), Inches(0.8))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_NAME
    run.font.size = title_size
    run.font.bold = True
    # Separator
    shape = slide.shapes.add_shape(1, Inches(0.7), Inches(1.1), Inches(12), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape.line.fill.background()
    # Bullets
    tf2 = add_textbox(slide, Inches(0.9), Inches(1.4), Inches(11.5), Inches(5.5))
    for i, bullet in enumerate(bullets):
        p = tf2.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        if isinstance(bullet, tuple):
            run.text = bullet[0]
            run.font.size = bullet_size
            run.font.name = FONT_NAME
            run.font.bold = bullet[1] if len(bullet) > 1 else False
        else:
            run.text = bullet
            run.font.size = bullet_size
            run.font.name = FONT_NAME
    return slide

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
tf = add_textbox(slide, Inches(1), Inches(1.5), Inches(11.33), Inches(1.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "AIDA"
run.font.name = FONT_NAME
run.font.size = Pt(44)
run.font.bold = True

tf2 = add_textbox(slide, Inches(1), Inches(2.8), Inches(11.33), Inches(1))
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "An Affective Intelligent Dialogue Agent Using the Adaptive\nEmotion-Aware Intervention Framework for Mental Health Support"
run2.font.name = FONT_NAME
run2.font.size = Pt(24)

tf3 = add_textbox(slide, Inches(1), Inches(4.8), Inches(11.33), Inches(2))
for line in [
    "Department of Computer Science",
    "Landmark University, Omu-Aran, Kwara State",
    "",
    "YOHANAN KOLAWOLA",
    "21CD008347",
]:
    set_para(tf3, line, size=Pt(18), bold=False, alignment=PP_ALIGN.CENTER)

tf4 = add_textbox(slide, Inches(1), Inches(6.8), Inches(5.5), Inches(0.6))
p4 = tf4.paragraphs[0]
run4 = p4.add_run()
run4.text = "Supervisor: Dr. Emmanuel Igbekele"
run4.font.name = FONT_NAME
run4.font.size = Pt(16)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — Outline
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Title
tf = add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
p = tf.paragraphs[0]
run = p.add_run()
run.text = "OUTLINE"
run.font.name = FONT_NAME
run.font.size = Pt(28)
run.font.bold = True

outline_items = [
    ("01", "Introduction"),
    ("02", "Statement of Problem"),
    ("03", "Aim and Objectives"),
    ("04", "Literature Review"),
    ("05", "Methodology"),
    ("06", "Results and Discussion"),
    ("07", "Conclusion"),
]

# Layout in two columns
left_items = outline_items[:4]
right_items = outline_items[4:]

for col_idx, items in enumerate([left_items, right_items]):
    base_x = Inches(1.5) if col_idx == 0 else Inches(7.5)
    for row_idx, (num, label) in enumerate(items):
        y = Inches(1.5 + row_idx * 1.3)
        # Number
        tf_num = add_textbox(slide, base_x, y, Inches(1), Inches(0.5))
        p_num = tf_num.paragraphs[0]
        run_n = p_num.add_run()
        run_n.text = num
        run_n.font.name = FONT_NAME
        run_n.font.size = Pt(24)
        run_n.font.bold = True
        # Label
        tf_lbl = add_textbox(slide, Inches(base_x + Inches(0.8)), y, Inches(4), Inches(0.5))
        p_lbl = tf_lbl.paragraphs[0]
        run_l = p_lbl.add_run()
        run_l.text = label
        run_l.font.name = FONT_NAME
        run_l.font.size = Pt(24)
        run_l.font.bold = True

# ══════════════════════════════════════════════════════════════════
# SLIDES 3-4 — Introduction
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "INTRODUCTION",
    [
        "Mental health disorders affect 1 in 8 people globally (WHO, 2022); in Nigeria, fewer than 10% of those affected access care",
        "AI-powered chatbots (Woebot, Wysa, Tess) have shown efficacy in reducing depression and anxiety symptoms through CBT-aligned conversation",
        "Existing chatbots detect current emotion but lack trajectory awareness — they cannot distinguish between transient distress and escalating crisis",
        "AIDA addresses this gap through the Adaptive Emotion-Aware Intervention Framework (AEIF), a 6-layer pipeline that integrates emotion detection, trend analysis, and strategy selection",
        "Target population: Nigerian young adults (18–35), with support for Nigerian Pidgin English and Nigeria-specific crisis resources",
    ],
    bullet_size=Pt(20)
)

add_content_slide(
    "INTRODUCTION (Cont'd)",
    [
        "The AEIF framework contributes a novel architecture: trajectory-aware strategy selection in mental health chatbots",
        "AEIF vs Baseline comparison demonstrates that strategy-guidance produces qualitatively superior empathetic alignment",
        "The system achieved 77.08% emotion accuracy, 100% crisis detection, 100% trend accuracy on a 58-case evaluation set",
        "All 31 unit tests pass across 5 test modules (emotion, crisis, trend, strategy, pipeline)",
    ],
    bullet_size=Pt(20)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — Statement of Problem
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "STATEMENT OF PROBLEM",
    [
        "Nigeria has a psychiatrist-to-population ratio of approximately 1:1 million (WHO), leaving millions without access to mental health support",
        "Existing mental health chatbots detect a user's current emotional state but do not model emotional trajectory over time",
        "Without trend awareness, chatbots cannot differentiate between a user who is persistently distressed vs. one who is improving",
        "Existing systems lack culturally adapted features for Nigerian users, including Pidgin English support and local crisis resources",
        "Strategy selection in existing systems is not driven by trajectory data, limiting response quality in multi-turn conversations",
    ],
    bullet_size=Pt(20)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Aim and Objectives
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "AIM AND OBJECTIVES",
    [
        ("Aim: To design, implement, and evaluate the AEIF framework — a trajectory-aware, culturally adapted conversational agent for preliminary mental health support among Nigerian young adults.", False),
        "",
        ("Objectives:", True),
        "1. Analyse existing mental health chatbot architectures to identify the gap in trajectory-aware emotional support",
        "2. Design the 6-layer AEIF framework integrating emotion classification, trend analysis, and strategy selection",
        "3. Implement AIDA as a FastAPI backend with a React frontend, including Nigerian Pidgin support and local crisis resources",
        "4. Evaluate the system using automated metrics (77.08% emotion accuracy, 100% crisis accuracy, 100% trend accuracy)",
        "5. Conduct a qualitative AEIF vs. Baseline pipeline comparison across three scripted conversation scenarios",
    ],
    bullet_size=Pt(18)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — Literature Review
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "LITERATURE REVIEW",
    [
        "Woebot (Fitzpatrick et al., 2017) — RCT showing significant depression reduction using CBT-based chatbot, but no emotion trajectory tracking",
        "Wysa (Inkster et al., 2018) — Empathy-driven conversational agent with demonstrated real-world impact, but limited to single-turn emotion detection",
        "Emotional Support Conversation framework (Liu et al., 2021) — Proposed multi-strategy approach for ESC, but no trajectory-to-strategy mapping",
        "GoEmotions dataset (Demszky et al., 2020) — 27 fine-grained emotion categories, basis for AEIF's 6-cluster taxonomy",
        "Key gap: No existing system integrates emotion trajectory analysis directly into strategy selection for mental health chatbots",
    ],
    bullet_size=Pt(18)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — Methodology: AEIF Framework
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "METHODOLOGY: AEIF FRAMEWORK",
    [
        "The AEIF framework is a 6-layer sequential processing pipeline:",
        "",
        "  Layer 1 — Safety Check: Keyword-based crisis detection (including Nigerian Pidgin)",
        "  Layer 2 — Emotion Detection: HF transformer + keyword fallback → 6 emotion clusters",
        "  Layer 3 — Memory Update: In-memory session store, 20-turn cap, dominant emotion tracking",
        "  Layer 4 — Trend Analysis: Severity-weighted scoring → 7 trend labels",
        "  Layer 5 — Strategy Selection: 12-condition lookup → 9 CBT-aligned strategies",
        "  Layer 6 — Response Generation: Groq LLM (llama-3.3-70b) + template fallback",
        "",
        "Design Science Research (DSR) methodology: 6 iterations of design-evaluate-refine",
    ],
    bullet_size=Pt(18)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — Emotion Clusters
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "METHODOLOGY: EMOTION TAXONOMY",
    [
        "AEIF defines 6 mutually exclusive emotion clusters mapped from HuggingFace and GoEmotions labels:",
        "",
        "  POSITIVE (0.0) — joy, love, gratitude, optimism, relief",
        "  NEUTRAL (0.3) — neutral statements, factual content",
        "  AMBIGUOUS (0.4) — confusion, surprise, contradictory expressions",
        "  ANXIETY (0.6) — fear, nervousness, worry, dread",
        "  ANGER (0.7) — anger, annoyance, disgust, frustration",
        "  SADNESS (0.8) — sadness, grief, disappointment, hopelessness",
        "  CRISIS (1.0) — detected by safety layer before emotion classification",
        "",
        "Ambiguity detection: contrastive conjunctions (but, though, however) + mixed keywords → AMBIGUOUS",
    ],
    bullet_size=Pt(18)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — Trend Analysis
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "METHODOLOGY: TREND ANALYSIS",
    [
        "Unweighted arithmetic mean of severity scores over sliding window (max 10 entries)",
        "Minimum 3 turns required before trend is derived",
        "",
        "7 trend types (detected in priority order):",
        "  1. ELEVATED_ANGER — 2+ consecutive ANGER (highest priority, safety-relevant)",
        "  2. ESCALATING_DISTRESS — every score strictly increasing",
        "  3. IMPROVING — every score strictly decreasing + negative emotions in first half",
        "  4. FLUCTUATING — 1+ direction reversals between consecutive triples",
        "  5. PERSISTENT_DISTRESS — mean severity > 0.55",
        "  6. STABLE_POSITIVE — mean severity ≤ 0.2 + POSITIVE in last 3 entries",
        "  7. INSUFFICIENT_DATA — fewer than 3 turns or no trend detected",
    ],
    bullet_size=Pt(18)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — Strategy Selection
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "METHODOLOGY: STRATEGY SELECTION",
    [
        "12-condition lookup table: 6 trend-to-strategy + 6 emotion-to-strategy (for INSUFFICIENT_DATA)",
        "",
        "  Trend → Strategy:",
        "  PERSISTENT_DISTRESS → VALIDATION_AND_REFLECTION",
        "  ESCALATING_DISTRESS → GROUNDING",
        "  ELEVATED_ANGER → CALM_REFLECTION",
        "  IMPROVING → ENCOURAGEMENT",
        "  STABLE_POSITIVE → POSITIVE_REINFORCEMENT",
        "  FLUCTUATING → EXPLORATORY_INQUIRY",
        "",
        "LLM system prompt includes: emotion cluster, trend, strategy label, strategy instruction, and 6 hard rules",
    ],
    bullet_size=Pt(18)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — Implementation
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "IMPLEMENTATION HIGHLIGHTS",
    [
        "Backend: FastAPI (Python 3.14), 7 REST endpoints, CORS-configured for localhost:5173",
        "Frontend: React + Vite 5, 6 components (Disclaimer, EmotionPulse, TrendTimeline, StrategyBadge, FeedbackStars, SessionSummary)",
        "Emotion Classifier: HuggingFace j-hartmann/emotion-english-distilroberta-base + keyword fallback",
        "LLM: Groq API (llama-3.3-70b-versatile) with exponential backoff retry (3 attempts, 0.5s/1.0s/2.0s)",
        "Crisis Resources: Nigeria-specific (MANI: 08091116264, LASEMA: 08000432584, NEMA: 112)",
        "Pidgin Support: Keyword lists for sadness (\"nobody dey for me\"), anxiety (\"my mind dey scatter\"), anger (\"I dey vex\")",
        "Baseline Comparison Mode: use_aeif=False disables Layers 4-5, enabling controlled A/B testing",
        "All data is in-memory only — no database, no disk writes, no PII collection",
    ],
    bullet_size=Pt(18)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — Results: Emotion Accuracy
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "RESULTS: EMOTION CLASSIFICATION",
    [
        "Automated evaluation on 58-case test set (48 non-crisis, 10 crisis):",
        "",
        "  Overall emotion accuracy: 77.08% (37/48 non-crisis cases)",
        "  Crisis detection accuracy: 100% (10/10 + no false positives)",
        "  Trend analysis accuracy: 100% (7/7 scripted sequences)",
        "  All 31 unit tests pass across 5 test modules",
        "",
        "Per-cluster performance:",
        "  POSITIVE: Precision 1.000, Recall 0.889, F1 = 0.941 (n=9)",
        "  ANGER:    Precision 0.800, Recall 1.000, F1 = 0.889 (n=4)",
        "  SADNESS:  Precision 1.000, Recall 0.800, F1 = 0.889 (n=10)",
        "  ANXIETY:  Precision 0.800, Recall 0.800, F1 = 0.800 (n=5)",
        "  NEUTRAL:  Precision 0.500, Recall 0.900, F1 = 0.643 (n=10)",
        "  AMBIGUOUS: Precision 1.000, Recall 0.400, F1 = 0.571 (n=10)",
    ],
    bullet_size=Pt(16)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — Results: AEIF vs Baseline
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "RESULTS: AEIF vs BASELINE COMPARISON",
    [
        "Three scripted scenarios (sadness, anxiety, mixed-emotion) run through both pipelines",
        "",
        "Key finding — Mixed-emotion Turn 2: Same input (\"Everyone says I should be happy\"), same emotion (POSITIVE)",
        "  AEIF (POSITIVE_REINFORCEMENT): \"It's great that you're surrounded by supportive people...\" — affirming, celebratory",
        "  Baseline (no strategy): \"It can be really tough when others don't understand...\" — explores incongruence",
        "",
        "Key finding — Sadness Turn 3: Three consecutive SADNESS → PERSISTENT_DISTRESS → VALIDATION_AND_REFLECTION",
        "  AEIF response: Reflects accumulated distress, names the emotional experience without interpreting",
        "  Baseline: Subtly pathologises (\"this can be a sign that you're feeling drained\")",
        "",
        "Pattern across all scenarios: AEIF responses are strategy-specific, context-escalating, and empathetically calibrated",
    ],
    bullet_size=Pt(16)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 15 — Demo Conversation
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "RESULTS: DEMO CONVERSATION (7-TURN)",
    [
        "Turn 1: \"I have been feeling really down lately\" → SADNESS → VALIDATION",
        "Turn 2: \"I failed my final year project presentation\" → SADNESS → VALIDATION",
        "Turn 3: \"My supervisor was disappointed\" → ANGER → PERSISTENT_DISTRESS → VALIDATION_AND_REFLECTION",
        "Turn 4: \"I feel so worthless and tired of trying\" → SADNESS → FLUCTUATING → EXPLORATORY_INQUIRY",
        "Turn 5: \"Maybe things can get better\" → NEUTRAL → FLUCTUATING → EXPLORATORY_INQUIRY",
        "Turn 6: \"My friend called and it helped a little\" → NEUTRAL → FLUCTUATING → EXPLORATORY_INQUIRY",
        "Turn 7: \"I think I just needed to talk to someone\" → NEUTRAL → FLUCTUATING → EXPLORATORY_INQUIRY",
        "",
        "Significant escalation at Turn 3: Accumulated distress triggers strategy shift from VALIDATION to VALIDATION_AND_REFLECTION",
        "Dominant emotion: SADNESS | Final trend: FLUCTUATING | 14 total messages",
    ],
    bullet_size=Pt(16)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 16 — Implementation Challenges
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "IMPLEMENTATION CHALLENGES",
    [
        "HuggingFace Model Latency: Transformer model startup 45-60s → resolved via HF_HUB_OFFLINE=1 cache (under 2s)",
        "Trend Threshold Calibration: Initial ESCALATING_DISTRESS was too strict for real plateaus → reordered priority chain",
        "Ambiguity Detection: Contradictory statements classified as NEUTRAL due to high HF neutral logit → added ambiguity pre-check",
        "Crisis Pattern Coverage: Indirect expressions missed → expanded to 35+ patterns including Pidgin",
        "Baseline Mode: Required structural change to pipeline (use_aeif flag) with zero regression risk to primary path",
    ],
    bullet_size=Pt(18)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 17 — Conclusion
# ══════════════════════════════════════════════════════════════════
add_content_slide(
    "CONCLUSION",
    [
        "AIDA implements the AEIF framework — a trajectory-aware, culturally adapted conversational agent for mental health support",
        "AEIF introduces a novel architecture: trend analysis directly drives strategy selection in multi-turn conversation",
        "Quantitative: 77.08% emotion accuracy, 100% crisis accuracy, 100% trend accuracy, 31/31 unit tests pass",
        "Qualitative: AEIF responses are consistently more empathetically calibrated than baseline across 3 scenarios",
        "Culturally adapted for Nigerian users: Pidgin English support, local crisis resources (MANI, LASEMA, NEMA)",
        "Limitations: In-memory storage only, no user study yet, keyword classifier misses subtle positive signals",
        "Future work: User study with Nigerian participants, contextual emotion classification, database persistence",
    ],
    bullet_size=Pt(18)
)

# ══════════════════════════════════════════════════════════════════
# SLIDE 18 — References
# ══════════════════════════════════════════════════════════════════
refs_text = [
    "Fitzpatrick, K. K., Darcy, A., & Vierhile, M. (2017). Delivering cognitive behavior therapy to young adults with symptoms of depression and anxiety using a fully automated conversational agent (Woebot): A randomized controlled trial. JMIR Mental Health, 4(2), e19.",
    "Inkster, B., Sarda, S., & Subramanian, V. (2018). An empathy-driven, conversational artificial intelligence agent (Wysa) for digital mental well-being. JMIR mHealth and uHealth, 6(11), e12106.",
    "Liu, S., Zheng, C., Demasi, O., Sabour, S., Li, Y., Yu, Z., ... & Huang, M. (2021). Towards emotional support dialog systems. Proceedings of ACL 2021.",
    "Demszky, D., Movshovitz-Attias, D., Ko, J., Cowen, A., Nemade, G., & Ravi, S. (2020). GoEmotions: A dataset of fine-grained emotions. Proceedings of ACL 2020.",
    "Patel, V., Saxena, S., Lund, C., et al. (2018). The Lancet Commission on global mental health and sustainable development. The Lancet, 392(10157), 1553-1598.",
    "Naslund, J. A., Aschbrenner, K. A., Araya, R., et al. (2017). Digital technology for treating and preventing mental disorders in low-income and middle-income countries. The Lancet Psychiatry, 4(6), 486-500.",
]
slide = prs.slides.add_slide(blank_layout)
tf = add_textbox(slide, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6))
p = tf.paragraphs[0]
run = p.add_run()
run.text = "REFERENCES"
run.font.name = FONT_NAME
run.font.size = Pt(24)
run.font.bold = True

shape = slide.shapes.add_shape(1, Inches(0.7), Inches(0.9), Inches(12), Pt(2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
shape.line.fill.background()

tf2 = add_textbox(slide, Inches(0.9), Inches(1.2), Inches(11.5), Inches(5.5))
for ref in refs_text:
    p = tf2.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(10)
    run = p.add_run()
    run.text = ref
    run.font.size = Pt(14)
    run.font.name = FONT_NAME

# ══════════════════════════════════════════════════════════════════
# SLIDE 19 — Acknowledgement
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
tf = add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7))
run = tf.paragraphs[0].add_run()
run.text = "ACKNOWLEDGEMENT"
run.font.name = FONT_NAME
run.font.size = Pt(24)
run.font.bold = True

ack_lines = [
    "I want to sincerely thank God Almighty for His grace throughout this project.",
    "",
    "Landmark University, Department of Computer Science",
    "",
    "Supervisor: Dr. Emmanuel O. Igbekele",
    "Co-supervisor: Mr. Fortune O. Amrevuawho",
    "HOD: Prof. Marion O. Adebiyi",
    "",
    "And finally my family and friends for their support.",
]
tf2 = add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(5))
for line in ack_lines:
    p = tf2.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(6)
    run = p.add_run()
    run.text = line
    run.font.name = FONT_NAME
    run.font.size = Pt(20)

# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════
output_path = "aida/AIDA_Defense_Slides.pptx"
prs.save(output_path)
print(f"✓ Saved {len(prs.slides)} slides to {output_path}")
